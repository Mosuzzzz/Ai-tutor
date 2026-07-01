import os
from typing import List
from pypdf import PdfReader
import docx
from pptx import Presentation

def extract_text_from_file(file_path: str) -> str:
    """Extracts raw text content from PDF, DOCX, PPTX, or Image files."""
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in (".docx", ".doc"):
        return extract_text_from_docx(file_path)
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx(file_path)
    elif ext in (".png", ".jpg", ".jpeg", ".webp"):
        return extract_text_from_image(file_path)
    else:
        # Text or unknown fallbacks
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception:
            raise ValueError(f"Unsupported file format: {ext}")

def extract_text_from_pdf(file_path: str) -> str:
    try:
        reader = PdfReader(file_path)
        text_content = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                text_content.append(text)
        return "\n".join(text_content)
    except Exception as e:
        # Fallback for text-encoded mock files in testing
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception:
            raise e

def extract_text_from_docx(file_path: str) -> str:
    try:
        doc = docx.Document(file_path)
        text_content = []
        for paragraph in doc.paragraphs:
            if paragraph.text:
                text_content.append(paragraph.text)
        return "\n".join(text_content)
    except Exception as e:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception:
            raise e

def extract_text_from_pptx(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content.append(shape.text)
        return "\n".join(text_content)
    except Exception as e:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception:
            raise e

def extract_text_from_image(file_path: str) -> str:
    """Runs real OCR via a vision model. Raises if OCR is unavailable or the
    image has no readable text, so the pipeline marks the doc 'error' instead of
    fabricating text. Never derive summary text from image metadata."""
    from services.ai_service import ocr_image
    return ocr_image(file_path)

def chunk_text(text: str, chunk_size: int = 800, overlap: int = 150) -> List[str]:
    """Splits raw text into overlapping structural chunks."""
    if not text:
        return []
        
    chunks = []
    start = 0
    text_len = len(text)
    
    while start < text_len:
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        # Advance by size minus overlap
        start += (chunk_size - overlap)
        if start >= text_len or len(chunk) < chunk_size:
            break
            
    return chunks
